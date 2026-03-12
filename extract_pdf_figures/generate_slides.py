import fitz
import os
import re
from io import BytesIO
from PIL import Image

from pptx import Presentation
from pptx.util import Inches


# --------------------------------------------------
# Configuration
# --------------------------------------------------

INPUT_PDF = "Scientific_Reports.pdf"

OUTPUT_PPT = "paper_presentation.pptx"

FIGURE_FOLDER = "extracted_figures"

MIN_IMAGE_WIDTH = 200
MIN_IMAGE_HEIGHT = 200


# --------------------------------------------------
# Utilities
# --------------------------------------------------

def ensure_folder(path):
    if not os.path.exists(path):
        os.makedirs(path)


def clean_caption(text):

    text = text.replace("\n", " ")
    text = re.sub(r"\s+", " ", text)

    return text.strip()


def find_caption_near_image(page, img_rect):

    blocks = page.get_text("blocks")

    caption_candidates = []

    for b in blocks:

        x0, y0, x1, y1, text, *_ = b

        rect = fitz.Rect(x0, y0, x1, y1)

        distance = abs(rect.y0 - img_rect.y1)

        if distance < 120:

            if re.search(r"(fig|figure)\s*\d+", text, re.IGNORECASE):

                caption_candidates.append((distance, text))

    if caption_candidates:

        caption_candidates.sort(key=lambda x: x[0])

        return clean_caption(caption_candidates[0][1])

    return "Figure"


# --------------------------------------------------
# Figure Extraction
# --------------------------------------------------

def extract_figures(pdf_path):

    doc = fitz.open(pdf_path)

    extracted = []

    for page_index in range(len(doc)):

        page = doc.load_page(page_index)

        image_list = page.get_images(full=True)

        for img_index, img in enumerate(image_list):

            xref = img[0]

            base_image = doc.extract_image(xref)

            image_bytes = base_image["image"]

            image = Image.open(BytesIO(image_bytes))

            if image.width < MIN_IMAGE_WIDTH or image.height < MIN_IMAGE_HEIGHT:
                continue

            filename = f"page{page_index+1:03d}_fig{img_index+1}.png"

            filepath = os.path.join(FIGURE_FOLDER, filename)

            image.save(filepath)

            img_rect = page.get_image_rects(xref)[0]

            caption = find_caption_near_image(page, img_rect)

            extracted.append({
                "image": filepath,
                "caption": caption,
                "page": page_index + 1
            })

    return extracted


# --------------------------------------------------
# Slide Generation
# --------------------------------------------------

def create_presentation(figures):

    prs = Presentation()

    blank_slide = prs.slide_layouts[6]

    for i, fig in enumerate(figures):

        slide = prs.slides.add_slide(blank_slide)

        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.3),
            Inches(9),
            Inches(0.8)
        )

        title = title_box.text_frame
        title.text = f"Figure {i+1}"

        image_path = fig["image"]

        slide.shapes.add_picture(
            image_path,
            Inches(1),
            Inches(1.2),
            width=Inches(8)
        )

        caption_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(6),
            Inches(9),
            Inches(1.2)
        )

        caption_frame = caption_box.text_frame

        caption_frame.text = fig["caption"]

    prs.save(OUTPUT_PPT)


# --------------------------------------------------
# Main
# --------------------------------------------------

def main():

    ensure_folder(FIGURE_FOLDER)

    figures = extract_figures(INPUT_PDF)

    print(f"Figures extracted: {len(figures)}")

    create_presentation(figures)

    print("Presentation created:", OUTPUT_PPT)


if __name__ == "__main__":
    main()
